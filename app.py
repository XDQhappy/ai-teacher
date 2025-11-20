import streamlit as st
from datetime import date
from lesson_plan_generator import generate_lesson_plan_with_text_chunks
from socratic_ai import ask_socratic
from dcmee_ppt import generate_ppt_with_kimi
import pdfplumber
from docx import Document
from pptx import Presentation

# ======================
# Streamlit 页面配置
# ======================
st.set_page_config(page_title="AI 教学助手", layout="wide")

# ======================
# 页面状态控制
# ======================
if "page" not in st.session_state:
    st.session_state["page"] = "subject"

# ======================
# 文件文本提取函数
# ======================
def extract_text_from_file(uploaded_file):
    name = uploaded_file.name.lower()
    text = ""
    if name.endswith(".pdf"):
        with pdfplumber.open(uploaded_file) as pdf:
            text = "\n".join([page.extract_text() or "" for page in pdf.pages])
    elif name.endswith(".docx") or name.endswith(".doc"):
        doc = Document(uploaded_file)
        text = "\n".join([p.text for p in doc.paragraphs])
    elif name.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    elif name.endswith(".txt"):
        text = uploaded_file.read().decode("utf-8")
    return text

def split_text_into_chunks(text, chunk_size=3000):
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

# ======================
# 页面1：科目选择
# ======================
if st.session_state["page"] == "subject":
    st.title("📚 选择学科和年级")

    subjects = {
        "初中数学": ["七年级上","七年级下","八年级上","八年级下","九年级上","九年级下"],
        "初中物理": [],
        "初中化学": []
    }

    # 对应年级图片
    grade_images = {
        "七年级上": "images/7_1.png",
        "七年级下": "images/7_2.png",
        "八年级上": "images/8_1.png"
    }

    for subject, grades in subjects.items():
        st.subheader(subject)
        if grades:
            cols = st.columns(len(grades))
            for idx, grade in enumerate(grades):
                with cols[idx]:
                    # 显示真实图片
                    img_path = grade_images.get(grade, "https://via.placeholder.com/150")
                    st.image(img_path, width=150)
                    if st.button(f"{grade}", key=f"{subject}_{grade}"):
                        st.session_state["lesson_subject"] = subject
                        st.session_state["lesson_grade"] = grade
                        st.session_state["page"] = "ai_lesson"

# ======================
# 页面2：AI 教案 + 苏格拉底
# ======================
elif st.session_state["page"] == "ai_lesson":

    if st.button("⬅️ 返回选择科目/年级"):
        st.session_state["page"] = "subject"
        st.session_state.pop("lesson_text", None)

    st.title("📘 AI 教案与苏格拉底问答")

    # --- 竖向展开控制 ---
    # 初始化状态
    if "show_socratic" not in st.session_state:
        st.session_state["show_socratic"] = False

    # --- 竖向展开按钮 ---
    _, right_col = st.columns([0.9, 0.1])
    with right_col:
        if st.button("💬\n展开\n问答", key="toggle", help="点击展开或隐藏苏格拉底问答区"):
            st.session_state["show_socratic"] = not st.session_state["show_socratic"]

    show_socratic = st.session_state["show_socratic"]

    # --- 动态布局 ---
    if show_socratic:
        col1, col2 = st.columns([2, 1])   # 展开后左 2 / 右 1
    else:
        col1, = st.columns([1])           # 未展开时只有左边


    # 左侧：教案 + Kimi PPT
    with col1:
        st.header("生成教案")

        # 自动填充学科和年级
        subject = st.session_state.get("lesson_subject", "")
        grade = st.session_state.get("lesson_grade", "")

        lesson_title = st.text_input("课程名称", placeholder="例如：有理数", key="lesson_title_input")
        lesson_type = st.radio(
            "授课类型",
            ["新授课", "复习课","习题课"],
            index=0,
            key="lesson_type_input"
        )
        subject_input = st.text_input("学科", value=subject, key="subject_input")
        grade_input = st.text_input("年级", value=grade, key="grade_input")
        duration = st.number_input("课时数", min_value=1, max_value=10, key="duration_input")
        key_vocab = st.text_input("关键词汇，用逗号分隔", key="key_vocab_input")
        teaching_goals = st.text_area("教学目标（可选）", key="teaching_goals_input")
        teaching_focus = st.text_area("教学重点（可选）", key="teaching_focus_input")
        teaching_difficulties = st.text_area("教学难点（可选）", key="teaching_difficulties_input")
        uploaded_file = st.file_uploader(
            "📎 上传辅助材料（PDF / Word / PPT / TXT）",
            type=["pdf","docx","pptx","txt"]
        )

        if st.button("🚀 生成教案", key="generate_lesson_btn"):
            if not lesson_title or not subject_input:
                st.warning("请填写课程名称和学科")
            else:
                if uploaded_file:
                    text = extract_text_from_file(uploaded_file)
                    text_chunks = split_text_into_chunks(text)
                else:
                    text_chunks = []

                with st.spinner("🧠 AI 正在生成教案，请稍候..."):
                    lesson_text = generate_lesson_plan_with_text_chunks(
                        lesson_title, subject_input, grade_input, duration,
                        key_vocab, text_chunks,
                        teaching_goals, teaching_focus, teaching_difficulties,
                        lesson_type
                    )
                    st.session_state["lesson_text"] = lesson_text
                st.success("✅ 教案生成完成！")

        if "lesson_text" in st.session_state:
            st.subheader("📘 教学方案")
            st.markdown(st.session_state["lesson_text"])

            st.download_button(
                "📥 下载教案文本",
                data=st.session_state["lesson_text"],
                file_name=f"教学方案_{date.today()}.txt",
                mime="text/plain",
                key="download_lesson"
            )

            # Kimi PPT 生成按钮
            if st.button("🎨 生成 PPT", key="generate_ppt"):
                with st.spinner("正在生成 PPT…"):
                    ppt_url = generate_ppt_with_kimi(
                        lesson_title, st.session_state["lesson_text"]
                    )
                    st.session_state["ppt_url"] = ppt_url
                st.success("✅ PPT 已生成！")

            # 清空按钮（教案已生成）
            if st.button("🗑️ 清空教案与 PPT", key="clear_after"):
                st.session_state.clear()
                st.success("已清空，可重新生成。")

            # 下载 PPT
            if "ppt_url" in st.session_state:
                st.markdown(f"[📥 点击下载 PPT]({st.session_state['ppt_url']})", unsafe_allow_html=True)

    # 右侧：苏格拉底问答
    if show_socratic:
        with col2:
            st.header("苏格拉底式问答")

            if "chat_history" not in st.session_state:
                st.session_state["chat_history"] = []

            with st.form(key="socratic_form_main"):
                user_question = st.text_input("向AI提问", placeholder="在此输入问题...")
                ask_button = st.form_submit_button("💬 提问")
                if ask_button and user_question.strip():
                    with st.spinner("AI 正在思考..."):
                        answer = ask_socratic(
                            st.session_state.get("lesson_text", ""),
                            user_question
                        )
                    st.session_state["chat_history"].append((user_question, answer))

            if st.session_state["chat_history"]:
                st.subheader("对话记录")
                for q, a in reversed(st.session_state["chat_history"]):
                    st.markdown(f"**学生:** {q}")
                    st.markdown(f"**AI:** {a}")
                    st.markdown("---")